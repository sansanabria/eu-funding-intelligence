"""
Generate PDF versions of outreach_tracker.xlsx and pitch_deck.pptx
for GitHub visibility. Uses the business color palette.
"""
import os
import sys

BASE = os.path.dirname(os.path.abspath(__file__))
XLSX = os.path.join(BASE, 'outreach_tracker.xlsx')
PPTX = os.path.join(BASE, 'pitch_deck.pptx')

# ── Brand palette ──
NAVY     = (35, 60, 100)
STEEL    = (75, 100, 140)
COPPER   = (170, 125, 60)
SAGE     = (55, 125, 85)
ROSE     = (155, 80, 80)
TEAL     = (55, 120, 130)
TXT      = (30, 35, 50)
TXT_MED  = (90, 100, 120)
TXT_LT   = (140, 150, 170)
BG       = (255, 255, 255)
SURFACE  = (247, 248, 252)
SURF2    = (237, 240, 247)

def rgb01(c):
    return (c[0]/255, c[1]/255, c[2]/255)


# ================================================================
#  PART 1 — OUTREACH TRACKER (XLSX → PDF)
# ================================================================
def generate_outreach_pdf():
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import mm, cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Spacer, Paragraph, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    import openpyxl

    wb = openpyxl.load_workbook(XLSX)
    out_path = os.path.join(BASE, 'outreach_tracker.pdf')

    def rc(t): return colors.Color(*rgb01(t))

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('Title', fontName='Helvetica-Bold', fontSize=16,
                                  textColor=rc(NAVY), spaceAfter=2*mm)
    subtitle_style = ParagraphStyle('Sub', fontName='Helvetica', fontSize=9,
                                     textColor=rc(TXT_MED), spaceAfter=4*mm)
    cell_style = ParagraphStyle('Cell', fontName='Helvetica', fontSize=6.5,
                                 textColor=rc(TXT), leading=8)
    cell_bold = ParagraphStyle('CellB', fontName='Helvetica-Bold', fontSize=6.5,
                                textColor=rc(TXT), leading=8)
    cell_copper = ParagraphStyle('CellC', fontName='Helvetica-Bold', fontSize=6.5,
                                  textColor=rc(COPPER), leading=8)
    cell_sage = ParagraphStyle('CellS', fontName='Helvetica-Bold', fontSize=6.5,
                                textColor=rc(SAGE), leading=8)
    cell_teal = ParagraphStyle('CellT', fontName='Helvetica-Bold', fontSize=6.5,
                                textColor=rc(TEAL), leading=8)
    cell_header = ParagraphStyle('CellH', fontName='Helvetica-Bold', fontSize=6.5,
                                  textColor=colors.white, leading=8)

    doc = SimpleDocTemplate(out_path, pagesize=landscape(A4),
                            leftMargin=12*mm, rightMargin=12*mm,
                            topMargin=15*mm, bottomMargin=12*mm)

    elements = []

    # ── Sheet configs ──
    sheet_configs = [
        {
            'name': 'Outreach Tracker',
            'title': 'EIF Fund Managers  //  Outreach Tracker',
            'subtitle': 'Q1 2025  //  European Commission - Your Europe Portal',
            'header_row': 4,  # 0-indexed
            'data_start': 5,
            'col_widths': [18*mm, 52*mm, 24*mm, 22*mm, 24*mm, 16*mm,
                           14*mm, 14*mm, 14*mm, 28*mm, 22*mm, 30*mm, 18*mm, 28*mm],
            'copper_cols': [4, 5],  # EIF, SCORE
            'sage_cols': [6, 7, 8],  # SOLAR, H2, DESAL
            'teal_cols': [],
        },
        {
            'name': 'Contact Details',
            'title': 'Contact Details  //  Fund Managers',
            'subtitle': '',
            'header_row': 3,
            'data_start': 4,
            'col_widths': [55*mm, 28*mm, 55*mm, 50*mm, 48*mm, 30*mm, 55*mm],
            'copper_cols': [],
            'sage_cols': [],
            'teal_cols': [4],  # EMAIL
        },
        {
            'name': 'Opportunity Matrix',
            'title': 'Opportunity Matrix  //  Project Fit',
            'subtitle': '',
            'header_row': 3,
            'data_start': 4,
            'col_widths': [60*mm, 30*mm, 22*mm, 40*mm, 40*mm, 40*mm],
            'copper_cols': [2],  # SCORE
            'sage_cols': [3, 4, 5],  # SOLAR, HYDROGEN, DESALINATION
            'teal_cols': [],
        },
    ]

    for si, cfg in enumerate(sheet_configs):
        ws = wb[cfg['name']]

        # Title + copper line
        elements.append(Paragraph(cfg['title'], title_style))
        # Copper accent line via a thin table
        accent = Table([['']],  colWidths=[250*mm], rowHeights=[1.2*mm])
        accent.setStyle(TableStyle([('BACKGROUND', (0,0), (0,0), rc(COPPER)),
                                     ('LINEBELOW', (0,0), (0,0), 0, rc(COPPER))]))
        elements.append(accent)
        if cfg['subtitle']:
            elements.append(Spacer(1, 2*mm))
            elements.append(Paragraph(cfg['subtitle'], subtitle_style))
        elements.append(Spacer(1, 3*mm))

        # Read header
        header_cells = []
        for cell in list(ws.rows)[cfg['header_row']]:
            val = str(cell.value) if cell.value else ''
            header_cells.append(Paragraph(val, cell_header))

        # Read data rows
        data_rows = []
        for row in list(ws.rows)[cfg['data_start']:]:
            vals = []
            for ci, cell in enumerate(row):
                val = str(cell.value) if cell.value is not None else ''
                if val == 'None':
                    val = ''
                # Style certain columns
                if ci in cfg['copper_cols'] and val:
                    vals.append(Paragraph(val, cell_copper))
                elif ci in cfg['sage_cols'] and val.upper() in ('YES', 'Y', '1', '1.0'):
                    vals.append(Paragraph('YES', cell_sage))
                elif ci in cfg['sage_cols'] and val in ('-', '0', '0.0', ''):
                    vals.append(Paragraph('—', ParagraphStyle('dim', fontName='Helvetica',
                                          fontSize=6.5, textColor=rc(TXT_LT), leading=8)))
                elif ci in cfg['teal_cols'] and val:
                    vals.append(Paragraph(val, cell_teal))
                else:
                    vals.append(Paragraph(val, cell_style))
            data_rows.append(vals)

        table_data = [header_cells] + data_rows
        col_widths = cfg['col_widths'][:len(header_cells)]

        t = Table(table_data, colWidths=col_widths, repeatRows=1)

        # Table style
        style_cmds = [
            # Header
            ('BACKGROUND', (0, 0), (-1, 0), rc(NAVY)),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 6.5),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 3*mm),
            ('TOPPADDING', (0, 0), (-1, 0), 3*mm),
            # Data rows
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 6.5),
            ('TOPPADDING', (0, 1), (-1, -1), 2*mm),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 2*mm),
            ('LEFTPADDING', (0, 0), (-1, -1), 2*mm),
            ('RIGHTPADDING', (0, 0), (-1, -1), 2*mm),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            # Grid
            ('GRID', (0, 0), (-1, -1), 0.3, rc(SURF2)),
        ]

        # Alternating row fills
        for i in range(1, len(table_data)):
            bg = SURFACE if i % 2 == 0 else BG
            style_cmds.append(('BACKGROUND', (0, i), (-1, i), rc(bg)))

        t.setStyle(TableStyle(style_cmds))
        elements.append(t)

        if si < len(sheet_configs) - 1:
            elements.append(PageBreak())

    doc.build(elements)
    print(f'  Saved: {out_path}')


# ================================================================
#  PART 2 — PITCH DECK (PPTX → PDF)
# ================================================================
def generate_pitch_pdf():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from reportlab.lib.pagesizes import landscape
    from reportlab.lib.units import mm, cm, inch
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import Paragraph
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    from io import BytesIO
    from PIL import Image
    import tempfile

    prs = Presentation(PPTX)
    out_path = os.path.join(BASE, 'pitch_deck.pdf')

    slide_w = prs.slide_width.inches
    slide_h = prs.slide_height.inches
    page_w = slide_w * inch
    page_h = slide_h * inch

    c = canvas.Canvas(out_path, pagesize=(page_w, page_h))

    def pptx_color_to_rgb(color_obj):
        """Extract RGB from a pptx color object."""
        try:
            if color_obj and color_obj.rgb:
                r = color_obj.rgb[0:2]
                g = color_obj.rgb[2:4]
                b = color_obj.rgb[4:6]
                return (int(r, 16)/255, int(g, 16)/255, int(b, 16)/255)
        except:
            pass
        return None

    def emu_to_inch(emu):
        if emu is None:
            return 0
        return emu / 914400

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx > 0:
            c.showPage()

        # White background
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)

        # Process shapes sorted by z-order (back to front)
        for shape in slide.shapes:
            left = emu_to_inch(shape.left) * inch
            top = emu_to_inch(shape.top) * inch
            width = emu_to_inch(shape.width) * inch
            height = emu_to_inch(shape.height) * inch

            # Convert top from PPTX (top-down) to PDF (bottom-up)
            pdf_y = page_h - top - height

            # Draw filled rectangles
            if shape.shape_type is not None:
                try:
                    fill = shape.fill
                    if fill.type is not None:
                        try:
                            fg = fill.fore_color
                            rgb = pptx_color_to_rgb(fg)
                            if rgb:
                                c.setFillColorRGB(*rgb)
                                c.rect(left, pdf_y, width, height, fill=1, stroke=0)
                        except:
                            pass
                except:
                    pass

            # Draw images
            if shape.shape_type == 13:  # Picture
                try:
                    img_blob = shape.image.blob
                    img = Image.open(BytesIO(img_blob))
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                        img.save(tmp.name, 'PNG')
                        c.drawImage(tmp.name, left, pdf_y, width, height,
                                   preserveAspectRatio=True, mask='auto')
                        os.unlink(tmp.name)
                except Exception as e:
                    pass

            # Draw text
            if shape.has_text_frame:
                tf = shape.text_frame
                text_y = pdf_y + height  # Start from top of shape

                for para in tf.paragraphs:
                    full_text = ''
                    font_size = 10
                    font_name = 'Helvetica'
                    font_color = (0, 0, 0)
                    is_bold = False

                    for run in para.runs:
                        full_text += run.text
                        if run.font.size:
                            font_size = run.font.size.pt
                        if run.font.bold:
                            is_bold = True
                        rgb = pptx_color_to_rgb(run.font.color)
                        if rgb:
                            font_color = rgb

                    if not full_text.strip():
                        text_y -= font_size * 1.3
                        continue

                    if is_bold:
                        font_name = 'Helvetica-Bold'

                    c.setFont(font_name, min(font_size, 48))
                    c.setFillColorRGB(*font_color)

                    # Alignment
                    text_x = left + 3
                    if para.alignment and para.alignment == 1:  # CENTER
                        tw = c.stringWidth(full_text, font_name, min(font_size, 48))
                        text_x = left + (width - tw) / 2

                    text_y -= font_size * 1.3
                    if text_y > 0 and text_y < page_h:
                        c.drawString(text_x, text_y, full_text)

            # Draw tables
            if shape.has_table:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)

                # Calculate cell dimensions
                col_widths = [emu_to_inch(col.width) * inch for col in table.columns]
                row_heights = [emu_to_inch(row.height) * inch for row in table.rows]

                cell_y = pdf_y + height
                for ri, row in enumerate(table.rows):
                    cell_y -= row_heights[ri]
                    cell_x = left

                    for ci, cell in enumerate(row.cells):
                        cw = col_widths[ci]
                        ch = row_heights[ri]

                        # Cell fill
                        try:
                            fill = cell.fill
                            if fill.type is not None:
                                rgb = pptx_color_to_rgb(fill.fore_color)
                                if rgb:
                                    c.setFillColorRGB(*rgb)
                                    c.rect(cell_x, cell_y, cw, ch, fill=1, stroke=0)
                        except:
                            # Header row default
                            if ri == 0:
                                c.setFillColorRGB(*rgb01(NAVY))
                                c.rect(cell_x, cell_y, cw, ch, fill=1, stroke=0)

                        # Cell border
                        c.setStrokeColorRGB(*rgb01(SURF2))
                        c.setLineWidth(0.3)
                        c.rect(cell_x, cell_y, cw, ch, fill=0, stroke=1)

                        # Cell text
                        cell_text = cell.text.strip()
                        if cell_text:
                            fs = 7
                            fn = 'Helvetica-Bold' if ri == 0 else 'Helvetica'
                            fc = (1, 1, 1) if ri == 0 else rgb01(TXT)

                            # Try to get actual font color
                            for p in cell.text_frame.paragraphs:
                                for r in p.runs:
                                    rc2 = pptx_color_to_rgb(r.font.color)
                                    if rc2:
                                        fc = rc2
                                    if r.font.bold:
                                        fn = 'Helvetica-Bold'
                                    if r.font.size:
                                        fs = min(r.font.size.pt, 9)

                            c.setFont(fn, fs)
                            c.setFillColorRGB(*fc)

                            # Truncate if needed
                            max_w = cw - 4
                            while c.stringWidth(cell_text, fn, fs) > max_w and len(cell_text) > 3:
                                cell_text = cell_text[:-4] + '...'

                            ty = cell_y + (ch - fs) / 2
                            c.drawString(cell_x + 2, ty, cell_text)

                        cell_x += cw

        # Page number
        c.setFont('Helvetica', 7)
        c.setFillColorRGB(*rgb01(TXT_LT))
        c.drawRightString(page_w - 15*mm, 8*mm, f'Page {slide_idx + 1}')

    c.save()
    print(f'  Saved: {out_path}')


# ================================================================
#  MAIN
# ================================================================
if __name__ == '__main__':
    print('Generating outreach_tracker.pdf...')
    generate_outreach_pdf()

    print('Generating pitch_deck.pdf...')
    generate_pitch_pdf()

    print('\nAll PDFs generated.')
